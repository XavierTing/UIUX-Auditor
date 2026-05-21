#!/usr/bin/env python3
"""
Generate a combined comparative PowerPoint report from multiple audit directories.

Usage:
    python3 generate_combined_pptx.py <output-file> <audit-dir-1> [<audit-dir-2> ...]

Example:
    python3 generate_combined_pptx.py combined-report.pptx \
        audits/2026-03-09_precious-metals-trading \
        audits/2026-03-10_precious-metals-trading-audited
"""

import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SEVERITY_COLORS = {
    "critical": RGBColor(0xDC, 0x26, 0x26),
    "high": RGBColor(0xEA, 0x58, 0x0C),
    "medium": RGBColor(0xCA, 0x8A, 0x04),
    "low": RGBColor(0x16, 0xA3, 0x4A),
}

SEVERITY_BG = {
    "critical": RGBColor(0xFE, 0xF2, 0xF2),
    "high": RGBColor(0xFF, 0xF7, 0xED),
    "medium": RGBColor(0xFE, 0xFC, 0xE8),
    "low": RGBColor(0xF0, 0xFD, 0xF4),
}

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_200 = RGBColor(0xE5, 0xE7, 0xEB)
BLUE_700 = RGBColor(0x1E, 0x40, 0xAF)
GREEN_700 = RGBColor(0x15, 0x80, 0x3D)
GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
GREEN_BORDER = RGBColor(0xBB, 0xF7, 0xD0)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.adjustments[0] = 0.02
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False,
                 color=GRAY_700, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def parse_audit(audit_dir):
    """Parse a single audit directory into structured data."""
    audit_dir = Path(audit_dir)
    screenshots_dir = audit_dir / "screenshots"
    data = {"dir": str(audit_dir), "name": audit_dir.name}

    # Load findings
    findings_path = audit_dir / "findings.json"
    data["findings"] = []
    if findings_path.exists():
        with open(findings_path) as f:
            data["findings"] = json.load(f)

    # Load report markdown
    report_path = audit_dir / "report.md"
    data["report_md"] = ""
    if report_path.exists():
        with open(report_path) as f:
            data["report_md"] = f.read()

    # Extract metadata
    data["title"] = "UI/UX Audit Report"
    data["meta"] = {}
    for line in data["report_md"].split("\n")[:20]:
        if line.startswith("# "):
            data["title"] = line[2:].strip()
        m = re.match(r"\*\*(.+?):\*\*\s*(.+)", line)
        if m:
            data["meta"][m.group(1).strip()] = m.group(2).strip().strip("`")

    # Executive summary
    exec_summary = ""
    in_exec = False
    for line in data["report_md"].split("\n"):
        if "EXECUTIVE SUMMARY" in line:
            in_exec = True
            continue
        if in_exec and line.startswith("## ") and "EXECUTIVE" not in line:
            break
        if in_exec and line.strip() and not line.startswith("---"):
            exec_summary += line + " "
    data["exec_summary"] = exec_summary.strip().replace("**", "")

    # Score
    score_match = re.search(r"Score:\s*(\d+\.?\d*)\s*/\s*10", exec_summary)
    data["score"] = score_match.group(1) if score_match else "N/A"

    # Severity counts
    sev_counts = {}
    for f in data["findings"]:
        s = f.get("severity", "medium").lower()
        sev_counts[s] = sev_counts.get(s, 0) + 1
    data["sev_counts"] = sev_counts

    # Annotated screenshots
    data["screens"] = []
    if screenshots_dir.exists():
        for png in sorted(screenshots_dir.glob("*-annotated.png")):
            screen_name = png.stem.replace("-annotated", "")
            findings_json = screenshots_dir / f"{screen_name}-findings.json"
            screen_findings = []
            if findings_json.exists():
                with open(findings_json) as fj:
                    sf = json.load(fj)
                findings_map = {fi["number"]: fi for fi in data["findings"]}
                for item in sf:
                    full = findings_map.get(item["number"], {})
                    screen_findings.append({
                        **item,
                        "finding": full.get("finding", item.get("label", "")),
                        "recommendation": full.get("recommendation", ""),
                        "severity": item.get("severity", full.get("severity", "medium")),
                    })
            data["screens"].append({
                "name": screen_name.replace("-", " ").title(),
                "path": str(png),
                "findings": sorted(screen_findings,
                    key=lambda x: ["critical", "high", "medium", "low"].index(
                        x.get("severity", "low").lower())),
            })

    # Sections
    sections = {}
    current_section = None
    current_content = []
    for line in data["report_md"].split("\n"):
        if line.startswith("## "):
            if current_section:
                sections[current_section] = "\n".join(current_content)
            current_section = line[3:].strip()
            current_content = []
        elif current_section:
            current_content.append(line)
    if current_section:
        sections[current_section] = "\n".join(current_content)
    data["sections"] = sections

    return data


def build_combined_pptx(output_file, audit_dirs):
    audits = [parse_audit(d) for d in audit_dirs]
    is_comparative = len(audits) > 1
    product_name = audits[0]["title"].replace("UI/UX AUDIT REPORT: ", "")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    title_text = f"Combined UI/UX Audit: {product_name}" if is_comparative else product_name
    add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                 title_text, font_size=34, bold=True, color=WHITE)
    add_text_box(slide, Inches(1), Inches(2.4), Inches(11), Inches(0.5),
                 "Comparative UI/UX Audit Report" if is_comparative else "Comprehensive UI/UX Audit Report",
                 font_size=18, color=RGBColor(0x94, 0xA3, 0xB8))

    # Meta
    meta = audits[-1]["meta"]
    date_display = " vs ".join(a["meta"].get("Audit Date", "N/A") for a in audits) if is_comparative else meta.get("Audit Date", "N/A")
    meta_items = [
        ("Audit Dates" if is_comparative else "Audit Date", date_display),
        ("Platform", meta.get("Platform", "N/A")),
        ("Auditor Persona", meta.get("Auditor Persona", "N/A")),
        ("Audits Compared" if is_comparative else "URL", f"{len(audits)} audits" if is_comparative else meta.get("URL", "N/A")),
    ]
    y = Inches(3.5)
    for label, value in meta_items:
        add_text_box(slide, Inches(1), y, Inches(2), Inches(0.3),
                     label.upper(), font_size=9, color=RGBColor(0x64, 0x74, 0x8B))
        add_text_box(slide, Inches(1), y + Inches(0.25), Inches(5), Inches(0.3),
                     value, font_size=13, color=WHITE)
        y += Inches(0.6)

    # Score display on the right
    if is_comparative:
        for i, a in enumerate(audits):
            x_pos = Inches(8.5) + i * Inches(2.5)
            s = float(a["score"]) if a["score"] != "N/A" else 0
            score_color = SEVERITY_COLORS["critical"] if s < 5 else (SEVERITY_COLORS["high"] if s < 7 else (SEVERITY_COLORS["medium"] if s < 8 else SEVERITY_COLORS["low"]))
            add_text_box(slide, x_pos, Inches(2.8), Inches(2), Inches(1.2),
                         a["score"], font_size=56, bold=True, color=score_color,
                         alignment=PP_ALIGN.CENTER)
            suffix = "Original" if i == 0 else "Re-audit"
            add_text_box(slide, x_pos, Inches(4.0), Inches(2), Inches(0.4),
                         suffix, font_size=12, color=RGBColor(0x94, 0xA3, 0xB8),
                         alignment=PP_ALIGN.CENTER)
    else:
        s = float(audits[0]["score"]) if audits[0]["score"] != "N/A" else 0
        score_color = SEVERITY_COLORS["critical"] if s < 5 else (SEVERITY_COLORS["high"] if s < 7 else (SEVERITY_COLORS["medium"] if s < 8 else SEVERITY_COLORS["low"]))
        add_text_box(slide, Inches(9.5), Inches(3.0), Inches(2.5), Inches(1.5),
                     audits[0]["score"], font_size=72, bold=True, color=score_color,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(9.5), Inches(4.5), Inches(2.5), Inches(0.4),
                     "UX Health Score / 10", font_size=12, color=RGBColor(0x94, 0xA3, 0xB8),
                     alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Score Comparison + Improvement ──
    if is_comparative:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)

        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     "Score Improvement", font_size=28, bold=True, color=DARK_BG)

        # Score cards side by side
        for i, a in enumerate(audits):
            x = Inches(1.5) + i * Inches(5.5)
            card = add_shape(slide, x, Inches(1.4), Inches(4), Inches(2.5),
                             fill_color=RGBColor(0xF9, 0xFA, 0xFB),
                             border_color=GRAY_200)
            s = float(a["score"]) if a["score"] != "N/A" else 0
            score_color = SEVERITY_COLORS["critical"] if s < 5 else (SEVERITY_COLORS["high"] if s < 7 else (SEVERITY_COLORS["medium"] if s < 8 else SEVERITY_COLORS["low"]))
            add_text_box(slide, x + Inches(0.5), Inches(1.6), Inches(3), Inches(1.2),
                         a["score"], font_size=64, bold=True, color=score_color,
                         alignment=PP_ALIGN.CENTER)
            suffix = "Original" if i == 0 else "Re-audit"
            add_text_box(slide, x + Inches(0.5), Inches(2.8), Inches(3), Inches(0.4),
                         f"{a['meta'].get('Audit Date', '')} ({suffix})",
                         font_size=14, color=GRAY_500, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, x + Inches(0.5), Inches(3.2), Inches(3), Inches(0.4),
                         f"{len(a['findings'])} findings",
                         font_size=12, color=GRAY_500, alignment=PP_ALIGN.CENTER)

        # Arrow between
        add_text_box(slide, Inches(5.8), Inches(2.0), Inches(1.5), Inches(1),
                     "→", font_size=48, bold=True, color=GREEN_700,
                     alignment=PP_ALIGN.CENTER)

        # Improvement banner
        try:
            delta = float(audits[-1]["score"]) - float(audits[0]["score"])
        except (ValueError, TypeError):
            delta = 0

        banner = add_shape(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5),
                           fill_color=GREEN_BG, border_color=GREEN_BORDER)
        add_text_box(slide, Inches(2), Inches(4.7), Inches(2), Inches(0.8),
                     f"+{delta:.1f}", font_size=42, bold=True, color=GREEN_700)
        add_text_box(slide, Inches(4), Inches(4.7), Inches(7), Inches(0.4),
                     "Score Improvement", font_size=18, bold=True, color=GREEN_700)

        total_orig = len(audits[0]["findings"])
        total_new = len(audits[-1]["findings"])
        add_text_box(slide, Inches(4), Inches(5.2), Inches(7), Inches(0.6),
                     f"Original: {total_orig} findings ({audits[0]['sev_counts'].get('critical', 0)} critical, {audits[0]['sev_counts'].get('high', 0)} high)  →  Re-audit: {total_new} findings ({audits[-1]['sev_counts'].get('critical', 0)} critical, {audits[-1]['sev_counts'].get('high', 0)} high)",
                     font_size=13, color=GRAY_700)

    # ── Slide 3: Severity Comparison ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Severity Breakdown" + (" Comparison" if is_comparative else ""),
                 font_size=28, bold=True, color=DARK_BG)

    if is_comparative:
        # Side-by-side severity cards for each audit
        for ai, a in enumerate(audits):
            x_base = Inches(0.8) + ai * Inches(6.2)
            suffix = "Original" if ai == 0 else "Re-audit"
            add_text_box(slide, x_base, Inches(1.3), Inches(5.8), Inches(0.4),
                         f"{a['meta'].get('Audit Date', '')} ({suffix}) — {len(a['findings'])} findings",
                         font_size=14, bold=True, color=GRAY_700)

            for si, (sev, label) in enumerate([("critical", "Critical"), ("high", "High"),
                                                ("medium", "Medium"), ("low", "Low")]):
                x = x_base + si * Inches(1.4)
                card = add_shape(slide, x, Inches(1.9), Inches(1.3), Inches(1.5),
                                 fill_color=SEVERITY_BG.get(sev),
                                 border_color=SEVERITY_COLORS.get(sev))
                add_text_box(slide, x + Inches(0.1), Inches(2.05), Inches(1.1), Inches(0.8),
                             str(a["sev_counts"].get(sev, 0)), font_size=36, bold=True,
                             color=SEVERITY_COLORS[sev], alignment=PP_ALIGN.CENTER)
                add_text_box(slide, x + Inches(0.1), Inches(2.8), Inches(1.1), Inches(0.3),
                             label.upper(), font_size=10, bold=True,
                             color=SEVERITY_COLORS[sev], alignment=PP_ALIGN.CENTER)
    else:
        x_start = Inches(0.8)
        card_width = Inches(2.8)
        card_gap = Inches(0.3)
        for i, (sev, label) in enumerate([("critical", "Critical"), ("high", "High"),
                                            ("medium", "Medium"), ("low", "Low")]):
            x = x_start + i * (card_width + card_gap)
            card = add_shape(slide, x, Inches(2.0), card_width, Inches(2.0),
                             fill_color=SEVERITY_BG.get(sev),
                             border_color=SEVERITY_COLORS.get(sev))
            add_text_box(slide, x + Inches(0.3), Inches(2.2), Inches(2.2), Inches(1.0),
                         str(audits[0]["sev_counts"].get(sev, 0)), font_size=48, bold=True,
                         color=SEVERITY_COLORS[sev], alignment=PP_ALIGN.CENTER)
            add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(2.2), Inches(0.4),
                         label.upper(), font_size=12, bold=True,
                         color=SEVERITY_COLORS[sev], alignment=PP_ALIGN.CENTER)

    # ── Slide: Finding Resolution Status (if comparative) ──
    if is_comparative and len(audits) >= 2:
        latest = audits[-1]
        status_section = ""
        for key in latest["sections"]:
            if "ORIGINAL FINDINGS STATUS" in key:
                status_section = latest["sections"][key]
                break

        if status_section:
            table_rows = re.findall(
                r"\|\s*(\d+)\s*\|\s*(.+?)\s*\|\s*(\w[\w\s]*)\s*\|\s*(.+?)\s*\|",
                status_section)
            valid_rows = [(n, f, s, d) for n, f, s, d in table_rows
                          if not n.startswith("-") and n != "Original #"]

            if valid_rows:
                slide = prs.slides.add_slide(blank_layout)
                set_slide_bg(slide, WHITE)

                add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                             "Finding Resolution Status", font_size=28, bold=True, color=DARK_BG)
                add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
                             "Cross-reference of original findings against the audited version",
                             font_size=13, color=GRAY_500)

                # Status summary counts
                status_counts = {}
                for _, _, status, _ in valid_rows:
                    s = status.strip().upper()
                    status_counts[s] = status_counts.get(s, 0) + 1

                sx = Inches(0.8)
                for status_label, status_color in [("FIXED", SEVERITY_COLORS["low"]),
                                                     ("PARTIALLY FIXED", SEVERITY_COLORS["medium"]),
                                                     ("NOT FIXED", SEVERITY_COLORS["high"]),
                                                     ("WORSE", SEVERITY_COLORS["critical"])]:
                    count = status_counts.get(status_label, 0)
                    if count > 0:
                        badge = add_shape(slide, sx, Inches(1.5), Inches(2.2), Inches(0.5),
                                          fill_color=status_color)
                        badge.text_frame.paragraphs[0].text = f"{status_label}: {count}"
                        badge.text_frame.paragraphs[0].font.size = Pt(11)
                        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
                        badge.text_frame.paragraphs[0].font.bold = True
                        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        badge.text_frame.paragraphs[0].font.name = "Calibri"
                        sx += Inches(2.4)

                # Table header
                y = Inches(2.3)
                headers = [("#", 0.5), ("Original Finding", 5.0), ("Status", 1.5), ("Notes", 5.0)]
                x = Inches(0.8)
                for h_text, h_width in headers:
                    add_text_box(slide, x, y, Inches(h_width), Inches(0.35),
                                 h_text, font_size=9, bold=True, color=GRAY_500)
                    x += Inches(h_width)

                y += Inches(0.45)
                # Paginate: ~8 rows per slide
                rows_per_slide = 8
                for ri, (num, finding, status, notes) in enumerate(valid_rows):
                    if ri > 0 and ri % rows_per_slide == 0:
                        # New slide for overflow
                        slide = prs.slides.add_slide(blank_layout)
                        set_slide_bg(slide, WHITE)
                        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                                     "Finding Resolution Status (cont.)", font_size=28, bold=True, color=DARK_BG)
                        y = Inches(1.3)
                        x = Inches(0.8)
                        for h_text, h_width in headers:
                            add_text_box(slide, x, y, Inches(h_width), Inches(0.35),
                                         h_text, font_size=9, bold=True, color=GRAY_500)
                            x += Inches(h_width)
                        y += Inches(0.45)

                    status_clean = status.strip().upper()
                    status_color_map = {
                        "FIXED": SEVERITY_COLORS["low"],
                        "PARTIALLY FIXED": SEVERITY_COLORS["medium"],
                        "NOT FIXED": SEVERITY_COLORS["high"],
                        "WORSE": SEVERITY_COLORS["critical"],
                    }
                    sc = status_color_map.get(status_clean, GRAY_500)

                    # Row background
                    row_h = Inches(0.55)
                    row_bg = add_shape(slide, Inches(0.8), y, Inches(11.8), row_h,
                                       fill_color=RGBColor(0xF9, 0xFA, 0xFB))

                    x = Inches(0.9)
                    add_text_box(slide, x, y + Inches(0.05), Inches(0.4), Inches(0.3),
                                 f"#{num}", font_size=11, bold=True, color=DARK_BG)
                    x += Inches(0.5)
                    add_text_box(slide, x, y + Inches(0.05), Inches(4.8), Inches(0.45),
                                 finding.strip(), font_size=10, color=GRAY_700)
                    x += Inches(5.0)

                    badge = add_shape(slide, x, y + Inches(0.1), Inches(1.3), Inches(0.28),
                                      fill_color=sc)
                    badge.text_frame.paragraphs[0].text = status_clean
                    badge.text_frame.paragraphs[0].font.size = Pt(8)
                    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
                    badge.text_frame.paragraphs[0].font.bold = True
                    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    badge.text_frame.paragraphs[0].font.name = "Calibri"
                    x += Inches(1.5)

                    add_text_box(slide, x, y + Inches(0.05), Inches(4.8), Inches(0.45),
                                 notes.strip(), font_size=9, color=GRAY_500)
                    y += row_h + Inches(0.05)

    # ── Per-audit findings + screenshots slides ──
    for ai, a in enumerate(audits):
        suffix = " (Original)" if ai == 0 and is_comparative else (" (Re-audit)" if ai > 0 else "")
        label = a["meta"].get("Audit Date", f"Audit {ai+1}")

        # Executive summary slide
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     f"{label}{suffix} — Executive Summary", font_size=24, bold=True, color=DARK_BG)

        summary_clean = a["exec_summary"]
        summary_clean = re.sub(r"Overall.*?/\s*10", "", summary_clean).strip()
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(3.5),
                     summary_clean, font_size=13, color=GRAY_700)

        # Findings slides (5 per slide)
        severity_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
        sorted_findings = sorted(a["findings"],
            key=lambda x: severity_order.get(x.get("severity", "low").lower(), 4))

        findings_per_slide = 5
        for batch_idx in range(0, len(sorted_findings), findings_per_slide):
            batch = sorted_findings[batch_idx:batch_idx + findings_per_slide]
            slide = prs.slides.add_slide(blank_layout)
            set_slide_bg(slide, WHITE)

            page_num = batch_idx // findings_per_slide + 1
            total_pages = (len(sorted_findings) + findings_per_slide - 1) // findings_per_slide
            add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                         f"{label}{suffix} — Findings ({page_num}/{total_pages})",
                         font_size=22, bold=True, color=DARK_BG)

            # Table header
            y = Inches(1.2)
            header_h = Inches(0.4)
            col_headers = [("#", 0.5), ("Screen", 2.2), ("Severity", 1.0), ("Finding", 4.5), ("Recommendation", 3.8)]
            x = Inches(0.8)
            for h_text, h_width in col_headers:
                add_text_box(slide, x, y, Inches(h_width), header_h,
                             h_text, font_size=9, bold=True, color=GRAY_500)
                x += Inches(h_width)

            y += header_h + Inches(0.05)
            row_h = Inches(1.1)
            for f in batch:
                sev = f.get("severity", "medium").lower()
                # Accent bar
                accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                Inches(0.8), y, Inches(0.08), row_h)
                accent.fill.solid()
                accent.fill.fore_color.rgb = SEVERITY_COLORS.get(sev, GRAY_500)
                accent.line.fill.background()

                row_bg = add_shape(slide, Inches(0.88), y, Inches(11.5), row_h,
                                   fill_color=RGBColor(0xF9, 0xFA, 0xFB))

                col_x = Inches(1.0)
                add_text_box(slide, col_x, y + Inches(0.1), Inches(0.4), Inches(0.3),
                             f"#{f['number']}", font_size=12, bold=True, color=DARK_BG)
                col_x += Inches(0.5)
                add_text_box(slide, col_x, y + Inches(0.1), Inches(2.0), Inches(0.9),
                             f.get("screen", ""), font_size=10, color=GRAY_700)
                col_x += Inches(2.2)

                badge_shape = add_shape(slide, col_x, y + Inches(0.15), Inches(0.85), Inches(0.28),
                                        fill_color=SEVERITY_COLORS.get(sev))
                badge_shape.text_frame.paragraphs[0].text = sev.upper()
                badge_shape.text_frame.paragraphs[0].font.size = Pt(8)
                badge_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
                badge_shape.text_frame.paragraphs[0].font.bold = True
                badge_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                badge_shape.text_frame.paragraphs[0].font.name = "Calibri"

                col_x += Inches(1.0)
                add_text_box(slide, col_x, y + Inches(0.08), Inches(4.3), Inches(0.95),
                             f.get("finding", ""), font_size=11, color=DARK_BG)
                col_x += Inches(4.5)
                add_text_box(slide, col_x, y + Inches(0.08), Inches(3.6), Inches(0.95),
                             f.get("recommendation", ""), font_size=10, color=GRAY_500)
                y += row_h + Inches(0.08)

        # Screenshot slides
        for screen in a["screens"]:
            slide = prs.slides.add_slide(blank_layout)
            set_slide_bg(slide, WHITE)

            add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
                         f"{label}{suffix} — {screen['name']}", font_size=22, bold=True, color=DARK_BG)

            img_path = screen["path"]
            if os.path.exists(img_path):
                from PIL import Image as PILImage
                with PILImage.open(img_path) as im:
                    img_w, img_h = im.size

                max_img_w = Inches(7.5)
                max_img_h = Inches(6.2)
                aspect = img_w / img_h
                if aspect > (max_img_w / max_img_h):
                    w = max_img_w
                    h = int(w / aspect)
                else:
                    h = max_img_h
                    w = int(h * aspect)

                slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.0), w, h)

            # Findings sidebar
            fx = Inches(8.2)
            fy = Inches(1.0)
            add_text_box(slide, fx, fy, Inches(4.8), Inches(0.4),
                         "Findings on this screen:", font_size=12, bold=True, color=GRAY_700)
            fy += Inches(0.5)

            for sf in screen["findings"]:
                sev = sf.get("severity", "medium").lower()
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, fx, fy + Inches(0.05),
                                             Inches(0.22), Inches(0.22))
                dot.fill.solid()
                dot.fill.fore_color.rgb = SEVERITY_COLORS.get(sev, GRAY_500)
                dot.line.fill.background()
                dot.text_frame.paragraphs[0].text = str(sf["number"])
                dot.text_frame.paragraphs[0].font.size = Pt(7)
                dot.text_frame.paragraphs[0].font.color.rgb = WHITE
                dot.text_frame.paragraphs[0].font.bold = True
                dot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                add_text_box(slide, fx + Inches(0.3), fy, Inches(4.5), Inches(0.22),
                             f"[{sev.upper()}] {sf.get('finding', sf.get('label', ''))}",
                             font_size=9, color=SEVERITY_COLORS.get(sev, GRAY_700))
                add_text_box(slide, fx + Inches(0.3), fy + Inches(0.22), Inches(4.5), Inches(0.4),
                             f"→ {sf.get('recommendation', '')}",
                             font_size=8, color=GRAY_500)
                fy += Inches(0.7)

    # ── Top 5 Recommendations (from latest audit) ──
    latest = audits[-1]
    top5_section = latest["sections"].get("TOP 5 PRIORITY RECOMMENDATIONS", "")
    recs = re.split(r"###\s*\d+\.\s*", top5_section)
    recs = [r.strip() for r in recs[1:] if r.strip()]

    if recs:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     f"Top 5 Priority Recommendations{' (Current)' if is_comparative else ''}",
                     font_size=28, bold=True, color=DARK_BG)

        y = Inches(1.3)
        for idx, rec in enumerate(recs[:5]):
            lines = rec.strip().split("\n")
            rec_title = lines[0].strip()
            effort = ""
            body_parts = []
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                em = re.match(r"- \*\*(.+?):\*\*\s*(.*)", line)
                if em:
                    if "Effort" in em.group(1):
                        effort = em.group(2)
                    else:
                        body_parts.append(f"{em.group(1)}: {em.group(2)}")

            card_h = Inches(1.0)
            card = add_shape(slide, Inches(0.8), y, Inches(11.5), card_h,
                             fill_color=RGBColor(0xF0, 0xF4, 0xFF),
                             border_color=RGBColor(0xBF, 0xDB, 0xFE))

            num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.15),
                                                 Inches(0.4), Inches(0.4))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = BLUE_700
            num_circle.line.fill.background()
            num_circle.text_frame.paragraphs[0].text = str(idx + 1)
            num_circle.text_frame.paragraphs[0].font.size = Pt(14)
            num_circle.text_frame.paragraphs[0].font.color.rgb = WHITE
            num_circle.text_frame.paragraphs[0].font.bold = True
            num_circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            add_text_box(slide, Inches(1.6), y + Inches(0.1), Inches(8), Inches(0.35),
                         rec_title, font_size=14, bold=True, color=BLUE_700)

            body_text = " | ".join(body_parts) if body_parts else ""
            add_text_box(slide, Inches(1.6), y + Inches(0.45), Inches(8.5), Inches(0.5),
                         body_text, font_size=10, color=GRAY_700)

            if effort:
                effort_color = RGBColor(0x16, 0x65, 0x34) if "Quick" in effort else (
                    RGBColor(0x92, 0x40, 0x0E) if "Medium" in effort else RGBColor(0x99, 0x1B, 0x1B))
                add_text_box(slide, Inches(10.5), y + Inches(0.15), Inches(1.5), Inches(0.3),
                             effort, font_size=10, bold=True, color=effort_color,
                             alignment=PP_ALIGN.RIGHT)

            y += card_h + Inches(0.12)

    # ── What's Working Well ──
    ww_section = latest["sections"].get("WHAT'S WORKING WELL", "")
    ww_items = re.findall(r"\d+\.\s*\*\*(.+?)\*\*\s*(.+?)(?=\n\d+\.|\Z)", ww_section, re.DOTALL)

    if ww_items:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     "What's Working Well", font_size=28, bold=True, color=DARK_BG)

        y = Inches(1.3)
        for ww_title, ww_desc in ww_items:
            card = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.0),
                             fill_color=GREEN_BG, border_color=GREEN_BORDER)
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(0.8), y, Inches(0.08), Inches(1.0))
            accent.fill.solid()
            accent.fill.fore_color.rgb = SEVERITY_COLORS["low"]
            accent.line.fill.background()

            add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(10.8), Inches(0.3),
                         ww_title.strip(), font_size=14, bold=True,
                         color=RGBColor(0x16, 0x65, 0x34))
            add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(10.8), Inches(0.5),
                         ww_desc.strip(), font_size=11, color=GRAY_700)
            y += Inches(1.15)

    # ── End slide ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1),
                 "Thank You", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5),
                 f"Generated by UX Audit Skill  ·  {'Comparative Audit' if is_comparative else meta.get('Audit Date', '')}",
                 font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)

    prs.save(output_file)
    file_size = os.path.getsize(output_file) / (1024 * 1024)
    print(f"Combined PPTX report generated: {output_file} ({file_size:.1f} MB)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 generate_combined_pptx.py <output-file> <audit-dir-1> [<audit-dir-2> ...]")
        sys.exit(1)
    build_combined_pptx(sys.argv[1], sys.argv[2:])
