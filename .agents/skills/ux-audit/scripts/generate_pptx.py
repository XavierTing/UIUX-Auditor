#!/usr/bin/env python3
"""
Generate a PowerPoint audit report with embedded annotated screenshots.

Usage:
    python3 generate_pptx.py <audit-dir> <output-file>

Example:
    python3 generate_pptx.py audits/2026-03-09_precious-metals/ audits/2026-03-09_precious-metals/report.pptx

Expects:
    <audit-dir>/report.md
    <audit-dir>/findings.json
    <audit-dir>/screenshots/*-annotated.png
"""

import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


SEVERITY_COLORS = {
    "critical": RGBColor(0xDC, 0x26, 0x26),
    "high": RGBColor(0xEA, 0x58, 0x0C),
    "medium": RGBColor(0xCA, 0x8A, 0x04),
    "low": RGBColor(0x16, 0xA3, 0x4A),
}

SEVERITY_BG = {
    "critical": RGBColor(0xFE, 0xF2, 0xF2),
    "high": RGBColor(0xFF, 0xF7, 0xED),
    "medium": RGBColor(0xFE, 0xFC, 0xE8),
    "low": RGBColor(0xF0, 0xFD, 0xF4),
}

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_200 = RGBColor(0xE5, 0xE7, 0xEB)
BLUE_700 = RGBColor(0x1E, 0x40, 0xAF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    # Small corner radius
    shape.adjustments[0] = 0.02
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False,
                 color=GRAY_700, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def build_pptx(audit_dir, output_file):
    audit_dir = Path(audit_dir)
    screenshots_dir = audit_dir / "screenshots"

    # Load data
    findings = []
    findings_path = audit_dir / "findings.json"
    if findings_path.exists():
        with open(findings_path) as f:
            findings = json.load(f)

    report_md = ""
    report_path = audit_dir / "report.md"
    if report_path.exists():
        with open(report_path) as f:
            report_md = f.read()

    # Parse metadata
    title = "UI/UX Audit Report"
    meta = {}
    for line in report_md.split("\n")[:20]:
        if line.startswith("# "):
            title = line[2:].strip()
        m = re.match(r"\*\*(.+?):\*\*\s*(.+)", line)
        if m:
            meta[m.group(1).strip()] = m.group(2).strip().strip("`")

    # Parse executive summary
    exec_summary = ""
    in_exec = False
    for line in report_md.split("\n"):
        if "EXECUTIVE SUMMARY" in line:
            in_exec = True
            continue
        if in_exec and line.startswith("## ") and "EXECUTIVE" not in line:
            break
        if in_exec and line.strip() and not line.startswith("---"):
            exec_summary += line + " "

    score_match = re.search(r"Score:\s*(\d+\.?\d*)\s*/\s*10", exec_summary)
    score = score_match.group(1) if score_match else "N/A"

    # Parse sections
    sections = {}
    current_section = None
    current_content = []
    for line in report_md.split("\n"):
        if line.startswith("## "):
            if current_section:
                sections[current_section] = "\n".join(current_content)
            current_section = line[3:].strip()
            current_content = []
        elif current_section:
            current_content.append(line)
    if current_section:
        sections[current_section] = "\n".join(current_content)

    # Severity counts
    sev_counts = {}
    for f in findings:
        s = f.get("severity", "medium")
        sev_counts[s] = sev_counts.get(s, 0) + 1

    # Collect annotated screenshots
    annotated_screens = []
    if screenshots_dir.exists():
        for png in sorted(screenshots_dir.glob("*-annotated.png")):
            screen_name = png.stem.replace("-annotated", "")
            findings_json = screenshots_dir / f"{screen_name}-findings.json"
            screen_findings = []
            if findings_json.exists():
                with open(findings_json) as fj:
                    sf = json.load(fj)
                findings_map = {fi["number"]: fi for fi in findings}
                for item in sf:
                    full = findings_map.get(item["number"], {})
                    screen_findings.append({
                        **item,
                        "finding": full.get("finding", item.get("label", "")),
                        "recommendation": full.get("recommendation", ""),
                        "severity": item.get("severity", full.get("severity", "medium")),
                    })
            annotated_screens.append({
                "name": screen_name.replace("-", " ").title(),
                "path": str(png),
                "findings": sorted(screen_findings,
                    key=lambda x: ["critical", "high", "medium", "low"].index(
                        x.get("severity", "low"))),
            })

    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 title, font_size=36, bold=True, color=WHITE)
    add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.5),
                 "Comprehensive UI/UX Audit Report", font_size=18, color=RGBColor(0x94, 0xA3, 0xB8))

    # Meta info
    meta_items = [
        ("Audit Date", meta.get("Audit Date", "N/A")),
        ("Platform", meta.get("Platform", "N/A")),
        ("Auditor Persona", meta.get("Auditor Persona", "N/A")),
        ("URL", meta.get("URL", "N/A")),
    ]
    y = Inches(3.8)
    for label, value in meta_items:
        add_text_box(slide, Inches(1), y, Inches(2), Inches(0.3),
                     label.upper(), font_size=9, color=RGBColor(0x64, 0x74, 0x8B))
        add_text_box(slide, Inches(1), y + Inches(0.25), Inches(5), Inches(0.3),
                     value, font_size=13, color=WHITE)
        y += Inches(0.65)

    # Score circle on right
    score_box = add_text_box(slide, Inches(9.5), Inches(3.2), Inches(2.5), Inches(1.5),
                             score, font_size=72, bold=True, color=RGBColor(0xEA, 0x58, 0x0C),
                             alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.5), Inches(4.7), Inches(2.5), Inches(0.4),
                 "UX Health Score / 10", font_size=12, color=RGBColor(0x94, 0xA3, 0xB8),
                 alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary + Severity Breakdown ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Executive Summary", font_size=28, bold=True, color=DARK_BG)

    # Summary text
    summary_clean = exec_summary.strip()
    # Remove the score line from summary text to avoid duplication
    summary_clean = re.sub(r"\*\*Overall.*?\*\*.*?$", "", summary_clean, flags=re.MULTILINE).strip()
    summary_clean = summary_clean.replace("**", "")

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8), Inches(2.5),
                 summary_clean, font_size=13, color=GRAY_700)

    # Severity cards
    x_start = Inches(0.8)
    card_width = Inches(2.8)
    card_gap = Inches(0.3)
    y_cards = Inches(4.5)

    for i, (sev, label) in enumerate([("critical", "Critical"), ("high", "High"),
                                        ("medium", "Medium"), ("low", "Low")]):
        x = x_start + i * (card_width + card_gap)
        card = add_shape(slide, x, y_cards, card_width, Inches(1.8),
                         fill_color=SEVERITY_BG.get(sev), border_color=SEVERITY_COLORS.get(sev))

        add_text_box(slide, x + Inches(0.3), y_cards + Inches(0.2), Inches(2), Inches(0.8),
                     str(sev_counts.get(sev, 0)), font_size=42, bold=True,
                     color=SEVERITY_COLORS[sev], alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), y_cards + Inches(1.1), Inches(2), Inches(0.4),
                     label.upper(), font_size=11, bold=True,
                     color=SEVERITY_COLORS[sev], alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
                 f"Total Findings: {len(findings)}", font_size=14, bold=True, color=GRAY_500)

    # ── Slide 3+: Findings by severity ──
    severity_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
    sorted_findings = sorted(findings, key=lambda x: severity_order.get(x.get("severity", "low"), 4))

    # Group findings into slides of ~5 each
    findings_per_slide = 5
    for batch_idx in range(0, len(sorted_findings), findings_per_slide):
        batch = sorted_findings[batch_idx:batch_idx + findings_per_slide]
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)

        page_num = batch_idx // findings_per_slide + 1
        total_pages = (len(sorted_findings) + findings_per_slide - 1) // findings_per_slide
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                     f"Findings ({page_num}/{total_pages})", font_size=24, bold=True, color=DARK_BG)

        # Table header
        y = Inches(1.2)
        header_h = Inches(0.4)
        headers = [("#", 0.5), ("Screen", 2.2), ("Severity", 1.0), ("Finding", 4.5), ("Recommendation", 3.8)]
        x = Inches(0.8)
        for h_text, h_width in headers:
            add_text_box(slide, x, y, Inches(h_width), header_h,
                         h_text, font_size=9, bold=True, color=GRAY_500)
            x += Inches(h_width)

        # Finding rows
        y += header_h + Inches(0.05)
        row_h = Inches(1.1)
        for f in batch:
            sev = f.get("severity", "medium")
            # Severity accent bar
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(0.8), y, Inches(0.08), row_h)
            accent.fill.solid()
            accent.fill.fore_color.rgb = SEVERITY_COLORS.get(sev, GRAY_500)
            accent.line.fill.background()

            # Row background
            row_bg = add_shape(slide, Inches(0.88), y, Inches(11.5), row_h,
                               fill_color=RGBColor(0xF9, 0xFA, 0xFB))

            col_x = Inches(1.0)
            add_text_box(slide, col_x, y + Inches(0.1), Inches(0.4), Inches(0.3),
                         f"#{f['number']}", font_size=12, bold=True, color=DARK_BG)
            col_x += Inches(0.5)
            add_text_box(slide, col_x, y + Inches(0.1), Inches(2.0), Inches(0.9),
                         f.get("screen", ""), font_size=10, color=GRAY_700)
            col_x += Inches(2.2)

            # Severity badge
            badge_shape = add_shape(slide, col_x, y + Inches(0.15), Inches(0.85), Inches(0.28),
                                    fill_color=SEVERITY_COLORS.get(sev))
            badge_shape.text_frame.paragraphs[0].text = sev.upper()
            badge_shape.text_frame.paragraphs[0].font.size = Pt(8)
            badge_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
            badge_shape.text_frame.paragraphs[0].font.bold = True
            badge_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge_shape.text_frame.paragraphs[0].font.name = "Calibri"

            col_x += Inches(1.0)
            add_text_box(slide, col_x, y + Inches(0.08), Inches(4.3), Inches(0.95),
                         f.get("finding", ""), font_size=11, color=DARK_BG)
            col_x += Inches(4.5)
            add_text_box(slide, col_x, y + Inches(0.08), Inches(3.6), Inches(0.95),
                         f.get("recommendation", ""), font_size=10, color=GRAY_500)
            y += row_h + Inches(0.08)

    # ── Screenshot slides ──
    for screen in annotated_screens:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)

        add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
                     screen["name"], font_size=24, bold=True, color=DARK_BG)

        # Add screenshot image (left side, fit proportionally)
        img_path = screen["path"]
        if os.path.exists(img_path):
            from PIL import Image as PILImage
            with PILImage.open(img_path) as im:
                img_w, img_h = im.size

            max_img_w = Inches(7.5)
            max_img_h = Inches(6.2)
            aspect = img_w / img_h
            if aspect > (max_img_w / max_img_h):
                w = max_img_w
                h = int(w / aspect)
            else:
                h = max_img_h
                w = int(h * aspect)

            slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.0), w, h)

        # Findings list on right side
        fx = Inches(8.2)
        fy = Inches(1.0)
        add_text_box(slide, fx, fy, Inches(4.8), Inches(0.4),
                     "Findings on this screen:", font_size=12, bold=True, color=GRAY_700)
        fy += Inches(0.5)

        for sf in screen["findings"]:
            sev = sf.get("severity", "medium")
            # Severity dot
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, fx, fy + Inches(0.05),
                                         Inches(0.22), Inches(0.22))
            dot.fill.solid()
            dot.fill.fore_color.rgb = SEVERITY_COLORS.get(sev, GRAY_500)
            dot.line.fill.background()
            dot.text_frame.paragraphs[0].text = str(sf["number"])
            dot.text_frame.paragraphs[0].font.size = Pt(7)
            dot.text_frame.paragraphs[0].font.color.rgb = WHITE
            dot.text_frame.paragraphs[0].font.bold = True
            dot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            add_text_box(slide, fx + Inches(0.3), fy, Inches(4.5), Inches(0.22),
                         f"[{sev.upper()}] {sf.get('finding', sf.get('label', ''))}",
                         font_size=9, color=SEVERITY_COLORS.get(sev, GRAY_700))
            add_text_box(slide, fx + Inches(0.3), fy + Inches(0.22), Inches(4.5), Inches(0.4),
                         f"→ {sf.get('recommendation', '')}",
                         font_size=8, color=GRAY_500)
            fy += Inches(0.7)

    # ── Top 5 Recommendations slide ──
    top5_section = sections.get("TOP 5 PRIORITY RECOMMENDATIONS", "")
    recs = re.split(r"###\s*\d+\.\s*", top5_section)
    recs = [r.strip() for r in recs[1:] if r.strip()]

    if recs:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     "Top 5 Priority Recommendations", font_size=28, bold=True, color=DARK_BG)

        y = Inches(1.3)
        for idx, rec in enumerate(recs[:5]):
            lines = rec.strip().split("\n")
            rec_title = lines[0].strip()

            # Extract effort
            effort = ""
            body_parts = []
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                em = re.match(r"- \*\*(.+?):\*\*\s*(.*)", line)
                if em:
                    if "Effort" in em.group(1):
                        effort = em.group(2)
                    else:
                        body_parts.append(f"{em.group(1)}: {em.group(2)}")

            # Card
            card_h = Inches(1.0)
            card = add_shape(slide, Inches(0.8), y, Inches(11.5), card_h,
                             fill_color=RGBColor(0xF0, 0xF4, 0xFF),
                             border_color=RGBColor(0xBF, 0xDB, 0xFE))

            # Number circle
            num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.15),
                                                 Inches(0.4), Inches(0.4))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = BLUE_700
            num_circle.line.fill.background()
            num_circle.text_frame.paragraphs[0].text = str(idx + 1)
            num_circle.text_frame.paragraphs[0].font.size = Pt(14)
            num_circle.text_frame.paragraphs[0].font.color.rgb = WHITE
            num_circle.text_frame.paragraphs[0].font.bold = True
            num_circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            add_text_box(slide, Inches(1.6), y + Inches(0.1), Inches(8), Inches(0.35),
                         rec_title, font_size=14, bold=True, color=BLUE_700)

            body_text = " | ".join(body_parts) if body_parts else ""
            add_text_box(slide, Inches(1.6), y + Inches(0.45), Inches(8.5), Inches(0.5),
                         body_text, font_size=10, color=GRAY_700)

            if effort:
                effort_color = RGBColor(0x16, 0x65, 0x34) if "Quick" in effort else (
                    RGBColor(0x92, 0x40, 0x0E) if "Medium" in effort else RGBColor(0x99, 0x1B, 0x1B))
                add_text_box(slide, Inches(10.5), y + Inches(0.15), Inches(1.5), Inches(0.3),
                             effort, font_size=10, bold=True, color=effort_color,
                             alignment=PP_ALIGN.RIGHT)

            y += card_h + Inches(0.12)

    # ── Accessibility Summary slide ──
    a11y_section = sections.get("ACCESSIBILITY SUMMARY", "")
    a11y_rows = re.findall(r"\|\s*(.+?)\s*\|\s*(\w+)\s*\|\s*(\w[\w\s]*)\s*\|\s*(.+?)\s*\|", a11y_section)

    if a11y_rows:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     "Accessibility Summary", font_size=28, bold=True, color=DARK_BG)

        y = Inches(1.3)
        # Header row
        cols = [("WCAG Criterion", 3.5), ("Level", 0.8), ("Status", 1.2), ("Details", 6.5)]
        x = Inches(0.8)
        for col_text, col_w in cols:
            add_text_box(slide, x, y, Inches(col_w), Inches(0.35),
                         col_text, font_size=9, bold=True, color=GRAY_500)
            x += Inches(col_w)

        y += Inches(0.45)

        for criterion, level, status, details in a11y_rows:
            if criterion.startswith("---") or criterion.startswith("WCAG"):
                continue
            status_clean = status.strip().upper()
            status_color = {
                "FAIL": SEVERITY_COLORS["critical"],
                "MISSING": SEVERITY_COLORS["high"],
                "UNCLEAR": SEVERITY_COLORS["medium"],
                "PASS": SEVERITY_COLORS["low"],
                "LIKELY PASS": SEVERITY_COLORS["low"],
            }.get(status_clean, GRAY_500)

            x = Inches(0.8)
            add_text_box(slide, x, y, Inches(3.5), Inches(0.5),
                         criterion.strip(), font_size=10, color=DARK_BG)
            x += Inches(3.5)
            add_text_box(slide, x, y, Inches(0.8), Inches(0.3),
                         level.strip(), font_size=10, color=GRAY_700)
            x += Inches(0.8)

            badge = add_shape(slide, x, y + Inches(0.02), Inches(0.9), Inches(0.25),
                              fill_color=status_color)
            badge.text_frame.paragraphs[0].text = status_clean
            badge.text_frame.paragraphs[0].font.size = Pt(7)
            badge.text_frame.paragraphs[0].font.color.rgb = WHITE
            badge.text_frame.paragraphs[0].font.bold = True
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.paragraphs[0].font.name = "Calibri"
            x += Inches(1.2)

            add_text_box(slide, x, y, Inches(6.5), Inches(0.5),
                         details.strip(), font_size=9, color=GRAY_700)
            y += Inches(0.55)

        add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(11), Inches(0.4),
                     "Overall Accessibility Risk Level: HIGH — Three AA failures + missing Level A requirement.",
                     font_size=13, bold=True, color=SEVERITY_COLORS["high"])

    # ── What's Working Well slide ──
    ww_section = sections.get("WHAT'S WORKING WELL", "")
    ww_items = re.findall(r"\d+\.\s*\*\*(.+?)\*\*\s*(.+?)(?=\n\d+\.|\Z)", ww_section, re.DOTALL)

    if ww_items:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, WHITE)
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                     "What's Working Well", font_size=28, bold=True, color=DARK_BG)

        y = Inches(1.3)
        for ww_title, ww_desc in ww_items:
            card = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.0),
                             fill_color=RGBColor(0xF0, 0xFD, 0xF4),
                             border_color=RGBColor(0xBB, 0xF7, 0xD0))

            # Green accent bar
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(0.8), y, Inches(0.08), Inches(1.0))
            accent.fill.solid()
            accent.fill.fore_color.rgb = SEVERITY_COLORS["low"]
            accent.line.fill.background()

            add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(10.8), Inches(0.3),
                         ww_title.strip(), font_size=14, bold=True,
                         color=RGBColor(0x16, 0x65, 0x34))
            add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(10.8), Inches(0.5),
                         ww_desc.strip(), font_size=11, color=GRAY_700)
            y += Inches(1.15)

    # ── End slide ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1),
                 "Thank You", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5),
                 f"Generated by UX Audit Skill  ·  {meta.get('Audit Date', '')}",
                 font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)

    # Save
    prs.save(output_file)
    file_size = os.path.getsize(output_file) / (1024 * 1024)
    print(f"PPTX report generated: {output_file} ({file_size:.1f} MB)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python3 generate_pptx.py <audit-dir> <output-file>")
        sys.exit(1)
    build_pptx(sys.argv[1], sys.argv[2])
